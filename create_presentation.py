from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.italic = True
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # Check if it's a bullet point
        if line.strip().startswith(('•', '-', '✅', '🚫', '💰', '📉', '⚖️', '💸', '🥬', '❓', '🔒', '⛓️', '💳', '🚚', '📱', '🌍', '🔔', '⭐', '💬', '📊', '📈', '🎯', '☁️', '🔄', '🌐', '🔐', '🚜', '🌾', '🏪', '🌱', '🤝', '📚', '🏘️', '♻️', '🚛', '🗺️', '📧', '📄', '🌐', '📊')):
            p.text = line.strip()
            p.level = 0
            p.font.size = Pt(18)
        elif line.strip().startswith(('1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.')):
            p.text = line.strip()
            p.level = 0
            p.font.size = Pt(20)
            p.font.bold = True
        elif line.strip().startswith('**') and line.strip().endswith('**'):
            # Bold header
            p.text = line.strip().replace('**', '')
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(51, 51, 51)
        elif line.strip().startswith('###'):
            # Subheading
            p.text = line.strip().replace('###', '').strip()
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 204)
        else:
            p.text = line.strip()
            p.font.size = Pt(16)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "Farmer2Market", "Direct Farm-to-Consumer Marketplace Platform")

# Slide 2: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "### Current Challenges in Agricultural Supply Chain",
    "",
    "**For Farmers:**",
    "🚫 Middlemen reduce profit margins by 30-50%",
    "💰 Delayed payments and financial uncertainty",
    "📉 Limited market reach and visibility",
    "⚖️ Unfair pricing and exploitation",
    "",
    "**For Consumers:**",
    "💸 Higher prices due to multiple intermediaries",
    "🥬 Lack of product freshness and quality assurance",
    "❓ No transparency in sourcing and pricing",
    "🔒 Limited trust in product authenticity"
])

# Slide 3: Our Solution
add_content_slide(prs, "Our Solution - Farmer2Market", [
    "### A Comprehensive Digital Marketplace Platform",
    "",
    "**Vision:** Eliminate middlemen and create a direct, transparent, and secure marketplace",
    "",
    "**Key Objectives:**",
    "1. Empower Farmers - Direct market access and fair pricing",
    "2. Benefit Consumers - Fresh products at competitive prices",
    "3. Ensure Security - Blockchain-based escrow for payment protection",
    "4. Build Trust - Transparent transactions and delivery tracking",
    "5. Enable Growth - Scalable platform for agricultural commerce"
])

# Slide 4: Core Features
add_content_slide(prs, "What We Have Done - Core Features", [
    "### 1. Multi-Role User System",
    "• Buyer Portal - Browse products, place orders, track deliveries",
    "• Farmer Dashboard - List products, manage inventory, confirm orders",
    "• Delivery Partner System - Accept assignments, update delivery status",
    "• Admin Panel - Manage users, assign deliveries, monitor platform",
    "",
    "### 2. Product Management",
    "• Product listing with images, descriptions, and pricing",
    "• Inventory management and stock tracking",
    "• Category-based product organization",
    "",
    "### 3. Order Management System",
    "• Complete order lifecycle tracking",
    "• Multiple order statuses (Pending → Paid → Confirmed → Delivered)",
    "• Order history and detailed order views"
])

# Slide 5: Advanced Features
add_content_slide(prs, "Advanced Features", [
    "### 4. Blockchain Escrow System ⛓️",
    "• Smart contract-based payment escrow",
    "• Funds held securely until delivery confirmation",
    "• Automatic fund release to farmers",
    "• Refund mechanism for disputes",
    "• Immutable transaction records",
    "",
    "### 5. Payment Integration 💳",
    "• Razorpay payment gateway integration",
    "• Multiple payment methods (UPI, Cards, Net Banking)",
    "• Cash on Delivery (COD) option",
    "• Secure payment verification",
    "",
    "### 6. Delivery Tracking System 🚚",
    "• Real-time delivery partner assignment",
    "• GPS-based location tracking",
    "• Delivery status updates"
])

# Slide 6: Technology Stack
add_content_slide(prs, "Technology Stack", [
    "### Backend Technologies",
    "• Framework: Flask (Python)",
    "• Database: PostgreSQL with SQLAlchemy ORM",
    "• Authentication: Flask-Login",
    "• API: RESTful architecture",
    "",
    "### Blockchain Layer",
    "• Smart Contracts: Solidity (^0.8.0)",
    "• Blockchain: Ethereum (Ganache for development)",
    "• Web3 Integration: Web3.py",
    "• Contract: OrderEscrow.sol (281 lines)",
    "",
    "### Payment Integration",
    "• Gateway: Razorpay",
    "• Security: HMAC-SHA256 signature validation"
])

# Slide 7: Blockchain Escrow Flow
add_content_slide(prs, "Blockchain Escrow Flow", [
    "### Smart Contract Workflow",
    "",
    "**1. Order Creation (PAID)**",
    "Buyer places order → Payment via Razorpay → Funds locked in blockchain escrow",
    "",
    "**2. Farmer Confirmation (FARMER_CONFIRMED)**",
    "Farmer reviews order → Confirms availability → Smart contract updated",
    "",
    "**3. Delivery Assignment (OUT_FOR_DELIVERY)**",
    "Admin assigns delivery partner → Partner accepts → Blockchain records delivery start",
    "",
    "**4. Delivery Completion (DELIVERED)**",
    "Partner submits proof of delivery → Proof hash stored on-chain",
    "",
    "**5. Fund Release (FUNDS_RELEASED)**",
    "Admin verifies delivery → Smart contract releases funds → Payment transferred to farmer"
])

# Slide 8: Why We Did It
add_content_slide(prs, "Why We Did It - Motivation & Benefits", [
    "### Why This Solution?",
    "",
    "**1. Financial Security**",
    "• Blockchain escrow protects both farmers and buyers",
    "• No payment fraud or defaults",
    "• Automated fund release based on delivery proof",
    "",
    "**2. Trust Building**",
    "• Immutable blockchain records",
    "• Verifiable delivery proofs",
    "• Transparent pricing",
    "",
    "**3. Efficiency Gains**",
    "• Eliminate 2-3 intermediaries",
    "• Direct farmer-to-consumer connection",
    "• Faster payment settlements"
])

# Slide 9: Value Proposition
add_content_slide(prs, "Value Proposition", [
    "### Value for Stakeholders",
    "",
    "**For Farmers (30-50% Profit Increase)**",
    "✅ Direct selling without middlemen",
    "✅ Fair market prices",
    "✅ Guaranteed payments via escrow",
    "",
    "**For Consumers (15-25% Cost Savings)**",
    "✅ Lower prices (no intermediary markup)",
    "✅ Fresh products directly from farms",
    "✅ Transparent sourcing",
    "",
    "**For Delivery Partners**",
    "✅ Flexible earning opportunities",
    "✅ Digital order management",
    "✅ Performance tracking"
])

# Slide 10: Smart Contract Features
add_content_slide(prs, "Smart Contract Implementation", [
    "### OrderEscrow Smart Contract",
    "",
    "**Key Functions:**",
    "• createOrder(orderId, buyer, farmer) - Locks payment in escrow",
    "• markFarmerConfirmed(orderId) - Updates order state",
    "• markOutForDelivery(orderId) - Tracks delivery initiation",
    "• submitDeliveryProof(orderId, proofHash) - Records delivery evidence",
    "• releaseFunds(orderId) - Transfers payment to farmer",
    "• refundBuyer(orderId) - Returns payment in disputes",
    "",
    "**Security Features:**",
    "• Re-entrancy protection",
    "• Access control (admin-only functions)",
    "• State validation checks",
    "• Event logging for transparency"
])

# Slide 11: Payment Integration
add_content_slide(prs, "Payment Integration", [
    "### Razorpay Integration",
    "",
    "**Payment Flow:**",
    "1. Order Creation - Backend creates Razorpay order",
    "2. Payment Processing - User completes payment",
    "3. Verification - Backend verifies HMAC-SHA256 signature",
    "4. Blockchain Integration - Payment sent to smart contract",
    "",
    "**Security:**",
    "• Signature verification prevents tampering",
    "• Webhook support for async updates",
    "• Refund API for cancellations",
    "• PCI-DSS compliant gateway"
])

# Slide 12: Security Features
add_content_slide(prs, "Multi-Layer Security Architecture", [
    "### Security Implementation",
    "",
    "**1. Authentication & Authorization**",
    "• Secure password hashing (Werkzeug)",
    "• Session-based authentication (Flask-Login)",
    "• Role-based access control (RBAC)",
    "",
    "**2. Payment Security**",
    "• Razorpay signature verification",
    "• HTTPS for payment pages",
    "• Secure credential storage",
    "",
    "**3. Blockchain Security**",
    "• Smart contract access control",
    "• Admin-only critical functions",
    "• Re-entrancy attack prevention"
])

# Slide 13: Project Statistics
add_content_slide(prs, "Project Statistics", [
    "### Development Metrics",
    "",
    "**Codebase:**",
    "• Backend Code: 15+ Python files",
    "• Database Models: 13 models",
    "• API Routes: 8 route modules",
    "• Smart Contract: 281 lines (Solidity)",
    "• Templates: 30+ HTML templates",
    "• Test Files: 7 comprehensive test suites",
    "",
    "**Features:**",
    "• User Roles: 4 (Buyer, Farmer, Delivery, Admin)",
    "• Order States: 9 lifecycle stages",
    "• Payment Methods: Multiple (UPI, Cards, COD)",
    "• Blockchain States: 8 escrow states"
])

# Slide 14: Challenges & Solutions
add_content_slide(prs, "Challenges & Solutions", [
    "### Technical Challenges Overcome",
    "",
    "**1. Blockchain Integration**",
    "Challenge: Synchronizing database state with blockchain",
    "Solution: Dual-write pattern with transaction hash tracking",
    "",
    "**2. Payment Security**",
    "Challenge: Ensuring payment authenticity",
    "Solution: HMAC signature verification + blockchain escrow",
    "",
    "**3. Escrow Timing**",
    "Challenge: When to release funds to farmers",
    "Solution: Multi-stage approval (delivery proof + admin verification)",
    "",
    "**4. Multi-Role Access**",
    "Challenge: Different dashboards for different users",
    "Solution: Role-based routing with Flask-Login decorators"
])

# Slide 15: Future Enhancements
add_content_slide(prs, "Future Enhancements", [
    "### Roadmap for Expansion",
    "",
    "**Phase 1: Enhanced Features**",
    "🌍 Multi-language support (Hindi, Kannada, Tamil, etc.)",
    "📱 Mobile application (iOS & Android)",
    "🔔 Push notifications",
    "⭐ Rating and review system",
    "",
    "**Phase 2: Advanced Analytics**",
    "📊 Farmer sales analytics dashboard",
    "📈 Market price trends and predictions",
    "🎯 Personalized product recommendations",
    "",
    "**Phase 3: Scalability**",
    "☁️ Cloud deployment (AWS/Azure)",
    "🔄 Load balancing and caching",
    "🌐 CDN for product images"
])

# Slide 16: Business Impact
add_content_slide(prs, "Business Impact", [
    "### Market Potential & Scalability",
    "",
    "**Target Market:**",
    "🌾 Farmers: 146 million farmers in India",
    "🛒 Consumers: 1.4 billion potential customers",
    "🚚 Delivery Partners: Gig economy workers",
    "",
    "**Revenue Model:**",
    "• Commission on each transaction (2-5%)",
    "• Premium farmer subscriptions",
    "• Featured product listings",
    "• Delivery service fees",
    "",
    "**Competitive Advantages:**",
    "• Blockchain-based trust and security",
    "• Direct farmer connection (no middlemen)",
    "• Integrated escrow system"
])

# Slide 17: Social Impact
add_content_slide(prs, "Social Impact", [
    "### Transforming Agricultural Commerce",
    "",
    "**Economic Impact:**",
    "💰 Increase farmer income by 30-50%",
    "📉 Reduce consumer costs by 15-25%",
    "💼 Create delivery partner employment",
    "",
    "**Social Benefits:**",
    "🌱 Support sustainable farming practices",
    "🤝 Build trust between farmers and consumers",
    "📚 Promote digital literacy in rural areas",
    "",
    "**Environmental Impact:**",
    "♻️ Reduce food waste in supply chain",
    "🚛 Optimize delivery routes (lower emissions)",
    "🌍 Promote local sourcing"
])

# Slide 18: Learning Outcomes
add_content_slide(prs, "Learning Outcomes", [
    "### Skills & Knowledge Gained",
    "",
    "**Technical Skills:**",
    "✅ Full-stack web development (Flask + Frontend)",
    "✅ Blockchain development (Solidity + Web3.py)",
    "✅ Payment gateway integration",
    "✅ Database design and ORM (SQLAlchemy)",
    "✅ RESTful API development",
    "✅ Smart contract deployment and testing",
    "",
    "**Software Engineering:**",
    "✅ MVC architecture pattern",
    "✅ Service-oriented design",
    "✅ Role-based access control",
    "✅ Transaction management"
])

# Slide 19: Demonstration Highlights
add_content_slide(prs, "Demonstration Highlights", [
    "### Key Workflows to Demonstrate",
    "",
    "**1. User Registration & Login**",
    "• Farmer registration with approval workflow",
    "• Buyer instant registration",
    "",
    "**2. Product Management**",
    "• Farmer adds new product",
    "• Product listing with images",
    "",
    "**3. Order Placement**",
    "• Buyer browses products",
    "• Checkout with Razorpay payment",
    "",
    "**4. Blockchain Escrow**",
    "• Payment locked in smart contract",
    "• View transaction on Ganache"
])

# Slide 20: Conclusion
add_content_slide(prs, "Conclusion", [
    "### Project Summary",
    "",
    "**What We Built:**",
    "A comprehensive, blockchain-powered marketplace platform connecting farmers",
    "directly with consumers, featuring secure escrow payments, real-time delivery",
    "tracking, and multi-stakeholder management.",
    "",
    "**Key Achievements:**",
    "✅ Functional blockchain escrow system",
    "✅ Integrated payment gateway",
    "✅ Multi-role user management",
    "✅ Complete order lifecycle tracking",
    "✅ Delivery partner coordination",
    "✅ Comprehensive testing and documentation",
    "",
    "**Impact:**",
    "Empowering farmers, benefiting consumers, and transforming agricultural",
    "supply chains through technology."
])

# Slide 21: Thank You
add_title_slide(prs, "Thank You", "Questions & Discussion")

# Save presentation
output_path = r"C:\farmer_market\Farmer2Market_Final_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
